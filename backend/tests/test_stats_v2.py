from io import BytesIO
import os
import uuid
from datetime import date
from types import SimpleNamespace

import pytest
import psycopg2
from psycopg2.extras import execute_values
from pptx import Presentation
from pptx.util import Inches
from openpyxl import load_workbook
from sqlalchemy import event

from app.services.statistics.definitions import percent
from app.services.statistics.pptx import percentage_color
from app.services.statistics.filters import build_filter
from app.services.statistics.service import StatisticsService

SYNC_DB_URL = os.environ["DATABASE_URL"].replace("postgresql+asyncpg://", "postgresql://")


def _exec(sql, params=None):
    with psycopg2.connect(SYNC_DB_URL) as connection:
        with connection.cursor() as cursor:
            cursor.execute(sql, params or {})


def _exec_values(sql, values, template=None):
    with psycopg2.connect(SYNC_DB_URL) as connection:
        with connection.cursor() as cursor:
            execute_values(cursor, sql, values, template=template)


@pytest.mark.parametrize(
    ("numerator", "denominator", "expected"),
    [(0, 0, 0), (1, 2, 50), (2, 3, 67), (1, 8, 13), (100, 928, 11), (199, 200, 100)],
)
def test_percent_uses_half_up(numerator, denominator, expected):
    assert percent(numerator, denominator) == expected


@pytest.mark.parametrize(
    ("value", "expected"),
    [(0, "E06666"), (19, "E06666"), (20, "F4B183"), (39, "F4B183"),
     (40, "F9CB9C"), (59, "F9CB9C"), (60, "FFD966"), (74, "FFD966"),
     (75, "A9D18E"), (89, "A9D18E"), (90, "63BE7B"), (100, "63BE7B"),
     (None, "D9E2F3")],
)
def test_pptx_percentage_color_boundaries(value, expected):
    assert percentage_color(value) == expected


@pytest.mark.asyncio
async def test_stats_contract_and_pptx(client, admin_headers):
    params = {"date_from": "2026-08-18", "date_to": "2026-08-20"}
    dashboard = await client.get("/api/v1/stats/dashboard", params=params, headers=admin_headers)
    assert dashboard.status_code == 200, dashboard.text
    payload = dashboard.json()
    assert payload["methodology"] == "v2"
    assert payload["timezone"] == "Europe/Moscow"
    assert payload["period"] == {"date_from": "2026-08-18", "date_to": "2026-08-20"}
    assert payload["totals"]["inspections_total"] == (
        payload["totals"]["inspections_green"]
        + payload["totals"]["inspections_with_defects"]
    )
    assert payload["totals"]["issues_not_fixed"] == (
        payload["totals"]["issues_found"] - payload["totals"]["issues_closed"]
    )

    dynamics = await client.get("/api/v1/stats/dynamics", params=params, headers=admin_headers)
    assert dynamics.status_code == 200, dynamics.text
    assert [day["date"] for day in dynamics.json()["days"]] == [
        "2026-08-18", "2026-08-19", "2026-08-20"
    ]

    categories = await client.get("/api/v1/stats/categories", params=params, headers=admin_headers)
    assert categories.status_code == 200, categories.text
    assert len(categories.json()["categories"]) == 9
    category_reference = await client.get("/api/v1/issues/categories", headers=admin_headers)
    assert category_reference.status_code == 200, category_reference.text
    assert [row["name"] for row in category_reference.json()][-1] == "Прочее"

    pptx = await client.get("/api/v1/stats/shtab.pptx", params=params, headers=admin_headers)
    assert pptx.status_code == 200, pptx.text
    assert "shtab_2026-08-18_2026-08-20.pptx" in pptx.headers["content-disposition"]
    presentation = Presentation(BytesIO(pptx.content))
    assert len(presentation.slides) == 2
    assert presentation.slide_width == Inches(13.333)
    assert presentation.slide_height == Inches(7.5)
    first_slide_text = "\n".join(shape.text for shape in presentation.slides[0].shapes if hasattr(shape, "text"))
    second_slide_text = "\n".join(shape.text for shape in presentation.slides[1].shapes if hasattr(shape, "text"))
    assert "УПРАВЛЕНИЕ ЖИЛИЩНО-КОММУНАЛЬНОГО ХОЗЯЙСТВА" in first_slide_text
    assert "1.1" in first_slide_text
    assert "1.2" in second_slide_text
    district_table = next(shape.table for shape in presentation.slides[0].shapes if shape.has_table)
    assert len(district_table.rows) == len(payload["districts"]) + 2
    assert "Чистые площадки" in [cell.text for cell in district_table.rows[0].cells]
    assert "МСК (UTC+3)" in first_slide_text
    assert "Устранено из выявленных" in second_slide_text
    assert "Доля требующих устранения" in second_slide_text

    excel = await client.get("/api/v1/reports/export.xlsx", params=params, headers=admin_headers)
    assert excel.status_code == 200, excel.text
    workbook = load_workbook(BytesIO(excel.content), data_only=True)
    summary = workbook["Сводка по районам"]
    assert [cell.value for cell in summary[1]][:17] == [
        "Район", "Площадок", "Проверено", "Охват %", "Обходов",
        "Чистые площадки", "% чистых площадок", "Площадки с нарушениями",
        "% площадок с нарушениями", "Без нарушений", "С наруш.", "Выявлено",
        "На проверке", "Требует устранения", "Просрочено",
        "Устранено из выявленных", "Доля требующих устранения",
    ]
    excel_by_name = {summary.cell(row, 1).value: summary.cell(row, 2).value
                     for row in range(2, summary.max_row + 1)}
    assert excel_by_name == {
        row["district_name"]: row["total_sites"] for row in payload["districts"]
    } | {"ИТОГО": payload["totals"]["total_sites"]}


@pytest.mark.asyncio
async def test_stats_rejects_invalid_or_excessive_period(client, admin_headers):
    reversed_period = await client.get(
        "/api/v1/stats/dashboard",
        params={"date_from": "2026-08-20", "date_to": "2026-08-19"},
        headers=admin_headers,
    )
    assert reversed_period.status_code == 422
    excessive = await client.get(
        "/api/v1/stats/dynamics",
        params={"date_from": "2025-01-01", "date_to": "2026-08-20"},
        headers=admin_headers,
    )
    assert excessive.status_code == 422


@pytest.mark.asyncio
async def test_stats_all_time_is_an_explicit_unbounded_period(client, admin_headers):
    response = await client.get(
        "/api/v1/stats/dashboard", params={"all_time": "true"}, headers=admin_headers,
    )
    assert response.status_code == 200, response.text
    assert response.json()["period"]["date_from"] == "2026-06-01"


@pytest.mark.asyncio
async def test_dashboard_site_quality_uses_latest_completed_inspection(client, admin_headers):
    me = await client.get("/api/v1/auth/me", headers=admin_headers)
    user_id = me.json()["id"]
    district_id, courtyard_id = [str(uuid.uuid4()) for _ in range(2)]
    site_ids = [str(uuid.uuid4()) for _ in range(94)]
    latest_inspection_ids = [str(uuid.uuid4()) for _ in site_ids]
    historical_inspection_ids = [str(uuid.uuid4()) for _ in range(30)]
    issue_ids = [str(uuid.uuid4()) for _ in historical_inspection_ids]

    _exec(
        "INSERT INTO districts(id,name,code) VALUES (%(district_id)s,%(district_name)s,%(code)s);"
        "INSERT INTO courtyards(id,district_id,name) VALUES (%(courtyard_id)s,%(district_id)s,'Двор качества');",
        {"district_id": district_id, "district_name": f"Quality {district_id[:8]}",
         "code": district_id[:8], "courtyard_id": courtyard_id},
    )
    _exec_values(
        "INSERT INTO sites(id,courtyard_id,type,area_m2,geometry,is_active) VALUES %s",
        [(site_id, courtyard_id, "Детская площадка", 100) for site_id in site_ids],
        template="(%s, %s, %s, %s, ST_GeomFromText("
        "'POLYGON((41 55,41.01 55,41.01 55.01,41 55.01,41 55))',4326), true)",
    )
    _exec_values(
        "INSERT INTO inspections(id,site_id,inspector_id,type,status,created_at,completed_at) VALUES %s",
        [
            (inspection_id, site_id, user_id, "regular", "completed",
             "2026-08-19T08:00:00Z", "2026-08-19T09:00:00Z")
            for site_id, inspection_id in zip(site_ids, latest_inspection_ids)
        ],
    )
    _exec_values(
        "INSERT INTO inspections(id,site_id,inspector_id,type,status,created_at,completed_at) VALUES %s",
        [
            (inspection_id, site_id, user_id, "regular", "completed",
             "2026-08-19T07:00:00Z" if index == 0 else "2026-08-18T08:00:00Z",
             "2026-08-19T08:00:00Z" if index == 0 else "2026-08-18T09:00:00Z")
            for index, (site_id, inspection_id) in enumerate(zip(site_ids, historical_inspection_ids))
        ],
    )
    _exec_values(
        "INSERT INTO issues(id,inspection_id,site_id,category_id,title,status,created_by,created_at) VALUES "
        "%s",
        [
            (issue_id, inspection_id, site_id, "Старое замечание", "open", user_id,
             "2026-08-18T10:00:00Z")
            for site_id, inspection_id, issue_id in zip(site_ids, historical_inspection_ids, issue_ids)
        ],
        template="(%s, %s, %s, (SELECT id FROM issue_categories WHERE name='Прочее'), %s, %s, %s, %s)",
    )

    response = await client.get(
        "/api/v1/stats/dashboard",
        params={"date_from": "2026-08-19", "date_to": "2026-08-19", "district_id": district_id},
        headers=admin_headers,
    )

    assert response.status_code == 200, response.text
    row = response.json()["districts"][0]
    # The first site has an older defective DONE inspection in this same MSK
    # period and a later clean one, so events number 95 while site quality
    # includes the site exactly once as clean.
    assert row["inspections_total"] == 95
    assert row["sites_inspected"] == 94
    assert row["sites_latest_clean"] == 94
    assert row["sites_latest_with_defects"] == 0
    assert row["clean_sites_pct"] == 100
    assert row["defect_sites_pct"] == 0
    assert row["issues_requires_work_current"] == 30


@pytest.mark.asyncio
async def test_dashboard_site_quality_is_null_without_completed_inspections(client, admin_headers):
    district_id = str(uuid.uuid4())
    _exec(
        "INSERT INTO districts(id,name,code) VALUES (%(district_id)s,%(district_name)s,%(code)s)",
        {"district_id": district_id, "district_name": f"Empty {district_id[:8]}",
         "code": district_id[:8]},
    )

    response = await client.get(
        "/api/v1/stats/dashboard",
        params={"date_from": "2026-08-19", "date_to": "2026-08-19", "district_id": district_id},
        headers=admin_headers,
    )

    assert response.status_code == 200, response.text
    row = response.json()["districts"][0]
    assert row["clean_sites_pct"] is None
    assert row["defect_sites_pct"] is None


@pytest.mark.asyncio
async def test_dashboard_uses_period_end_issue_snapshot_and_period_quality(client, admin_headers):
    me = await client.get("/api/v1/auth/me", headers=admin_headers)
    user_id = me.json()["id"]
    district_id, empty_district_id, courtyard_id, site_id = [str(uuid.uuid4()) for _ in range(4)]
    july_inspection_id, august_inspection_id, july_closed_id, july_open_id, august_issue_id = [
        str(uuid.uuid4()) for _ in range(5)
    ]

    _exec(
        "INSERT INTO districts(id,name,code) VALUES (%(district_id)s,%(district_name)s,%(code)s);"
        "INSERT INTO districts(id,name,code) VALUES (%(empty_district_id)s,'Пустой район',%(empty_code)s);"
        "INSERT INTO courtyards(id,district_id,name) VALUES (%(courtyard_id)s,%(district_id)s,'Двор периода');"
        "INSERT INTO sites(id,courtyard_id,type,area_m2,geometry,is_active) VALUES "
        "(%(site_id)s,%(courtyard_id)s,'Детская площадка',100,ST_GeomFromText("
        "'POLYGON((42 55,42.01 55,42.01 55.01,42 55.01,42 55))',4326),true);"
        "INSERT INTO inspections(id,site_id,inspector_id,type,status,created_at,completed_at) VALUES "
        "(%(july_inspection_id)s,%(site_id)s,%(user_id)s,'regular','completed','2026-07-15T08:00:00Z','2026-07-15T09:00:00Z'),"
        "(%(august_inspection_id)s,%(site_id)s,%(user_id)s,'regular','completed','2026-08-15T08:00:00Z','2026-08-15T09:00:00Z');"
        "INSERT INTO issues(id,inspection_id,site_id,category_id,title,status,due_date,created_by,created_at) VALUES "
        "(%(july_closed_id)s,%(july_inspection_id)s,%(site_id)s,(SELECT id FROM issue_categories WHERE name='Прочее'),'Закрыто в июле','closed','2026-07-20',%(user_id)s,'2026-07-15T10:00:00Z'),"
        "(%(july_open_id)s,%(july_inspection_id)s,%(site_id)s,(SELECT id FROM issue_categories WHERE name='Прочее'),'Закрыто в августе','closed','2026-07-20',%(user_id)s,'2026-07-16T10:00:00Z'),"
        "(%(august_issue_id)s,%(july_inspection_id)s,%(site_id)s,(SELECT id FROM issue_categories WHERE name='Прочее'),'Создано в августе','closed','2026-08-20',%(user_id)s,'2026-08-15T10:00:00Z');"
        "INSERT INTO issue_status_history(issue_id,old_status,new_status,changed_by,created_at) VALUES "
        "(%(july_closed_id)s,'open','closed',%(user_id)s,'2026-07-20T09:00:00Z'),"
        "(%(july_open_id)s,'open','fixed',%(user_id)s,'2026-08-10T09:00:00Z'),"
        "(%(july_open_id)s,'fixed','closed',%(user_id)s,'2026-08-20T09:00:00Z'),"
        "(%(august_issue_id)s,'open','closed',%(user_id)s,'2026-08-20T09:00:00Z');",
        {
            "district_id": district_id, "district_name": f"Snapshot {district_id[:8]}",
            "code": district_id[:8], "empty_district_id": empty_district_id,
            "empty_code": empty_district_id[:8], "courtyard_id": courtyard_id, "site_id": site_id,
            "july_inspection_id": july_inspection_id, "august_inspection_id": august_inspection_id,
            "july_closed_id": july_closed_id, "july_open_id": july_open_id,
            "august_issue_id": august_issue_id, "user_id": user_id,
        },
    )

    july = await client.get(
        "/api/v1/stats/dashboard",
        params={"date_from": "2026-07-01", "date_to": "2026-07-31", "district_id": district_id},
        headers=admin_headers,
    )
    august = await client.get(
        "/api/v1/stats/dashboard",
        params={"date_from": "2026-08-01", "date_to": "2026-08-31", "district_id": district_id},
        headers=admin_headers,
    )
    empty = await client.get(
        "/api/v1/stats/dashboard",
        params={"date_from": "2026-07-01", "date_to": "2026-07-31", "district_id": empty_district_id},
        headers=admin_headers,
    )

    assert july.status_code == 200, july.text
    assert august.status_code == 200, august.text
    assert empty.status_code == 200, empty.text
    july_row = july.json()["districts"][0]
    august_row = august.json()["districts"][0]
    assert july_row["clean_sites_pct"] == 0
    assert august_row["clean_sites_pct"] == 100
    assert july_row["issues_found"] == 2
    assert july_row["issues_cohort_closed_as_of"] == 1
    assert july_row["issues_cohort_closed_pct"] == 50
    assert july_row["issues_snapshot_total"] == 2
    assert july_row["issues_requires_work_pct"] == 50
    empty_row = empty.json()["districts"][0]
    assert empty_row["issues_cohort_closed_pct"] is None
    assert empty_row["issues_requires_work_pct"] is None
    # Today the Issue row is closed, but in July it was still initially open:
    # the period snapshot must read history as of the selected end, not Issue.status.
    assert july_row["issues_requires_work_current"] == 1
    assert july_row["issues_pending_final_current"] == 0
    assert july_row["issues_overdue_current"] == 1
    assert august_row["issues_requires_work_current"] == 0


@pytest.mark.asyncio
async def test_dashboard_remediation_percentages_use_snapshot_and_aggregate_denominators(
    client, admin_headers,
):
    me = await client.get("/api/v1/auth/me", headers=admin_headers)
    user_id = me.json()["id"]
    district_a_id, district_b_id, courtyard_a_id, courtyard_b_id, site_a_id, site_b_id = [
        str(uuid.uuid4()) for _ in range(6)
    ]
    inspection_a_id, inspection_b_id, old_open_id, july_closed_a_id, july_open_id, july_closed_b_id = [
        str(uuid.uuid4()) for _ in range(6)
    ]

    _exec(
        "INSERT INTO districts(id,name,code) VALUES "
        "(%(district_a_id)s,%(district_a_name)s,%(district_a_code)s),"
        "(%(district_b_id)s,%(district_b_name)s,%(district_b_code)s);"
        "INSERT INTO courtyards(id,district_id,name) VALUES "
        "(%(courtyard_a_id)s,%(district_a_id)s,'Двор A'),"
        "(%(courtyard_b_id)s,%(district_b_id)s,'Двор B');"
        "INSERT INTO sites(id,courtyard_id,type,area_m2,geometry,is_active) VALUES "
        "(%(site_a_id)s,%(courtyard_a_id)s,'Детская площадка',100,ST_GeomFromText("
        "'POLYGON((43 55,43.01 55,43.01 55.01,43 55.01,43 55))',4326),true),"
        "(%(site_b_id)s,%(courtyard_b_id)s,'Детская площадка',100,ST_GeomFromText("
        "'POLYGON((44 55,44.01 55,44.01 55.01,44 55.01,44 55))',4326),true);"
        "INSERT INTO inspections(id,site_id,inspector_id,type,status,created_at,completed_at) VALUES "
        "(%(inspection_a_id)s,%(site_a_id)s,%(user_id)s,'regular','completed','2000-07-15T08:00:00Z','2000-07-15T09:00:00Z'),"
        "(%(inspection_b_id)s,%(site_b_id)s,%(user_id)s,'regular','completed','2000-07-15T08:00:00Z','2000-07-15T09:00:00Z');"
        "INSERT INTO issues(id,inspection_id,site_id,category_id,title,status,created_by,created_at) VALUES "
        "(%(old_open_id)s,%(inspection_a_id)s,%(site_a_id)s,(SELECT id FROM issue_categories WHERE name='Прочее'),'Старое открытое','closed',%(user_id)s,'2000-06-15T10:00:00Z'),"
        "(%(july_closed_a_id)s,%(inspection_a_id)s,%(site_a_id)s,(SELECT id FROM issue_categories WHERE name='Прочее'),'Июльское закрытое A','closed',%(user_id)s,'2000-07-15T10:00:00Z'),"
        "(%(july_open_id)s,%(inspection_a_id)s,%(site_a_id)s,(SELECT id FROM issue_categories WHERE name='Прочее'),'Июльское открытое','closed',%(user_id)s,'2000-07-16T10:00:00Z'),"
        "(%(july_closed_b_id)s,%(inspection_b_id)s,%(site_b_id)s,(SELECT id FROM issue_categories WHERE name='Прочее'),'Июльское закрытое B','closed',%(user_id)s,'2000-07-15T10:00:00Z');"
        "INSERT INTO issue_status_history(issue_id,old_status,new_status,changed_by,created_at) VALUES "
        "(%(july_closed_a_id)s,'open','closed',%(user_id)s,'2000-07-20T10:00:00Z'),"
        "(%(july_closed_b_id)s,'open','closed',%(user_id)s,'2000-07-20T10:00:00Z');",
        {
            "district_a_id": district_a_id, "district_a_name": f"Aggregate A {district_a_id[:8]}",
            "district_a_code": district_a_id[:8], "district_b_id": district_b_id,
            "district_b_name": f"Aggregate B {district_b_id[:8]}", "district_b_code": district_b_id[:8],
            "courtyard_a_id": courtyard_a_id, "courtyard_b_id": courtyard_b_id,
            "site_a_id": site_a_id, "site_b_id": site_b_id,
            "inspection_a_id": inspection_a_id, "inspection_b_id": inspection_b_id,
            "old_open_id": old_open_id, "july_closed_a_id": july_closed_a_id,
            "july_open_id": july_open_id, "july_closed_b_id": july_closed_b_id,
            "user_id": user_id,
        },
    )

    response = await client.get(
        "/api/v1/stats/dashboard",
        params={"date_from": "2000-07-01", "date_to": "2000-07-31"},
        headers=admin_headers,
    )

    assert response.status_code == 200, response.text
    payload = response.json()
    rows = {row["district_id"]: row for row in payload["districts"]}
    district_a = rows[district_a_id]
    district_b = rows[district_b_id]
    assert district_a["issues_found"] == 2
    assert district_a["issues_snapshot_total"] == 3
    assert district_a["issues_requires_work_current"] == 2
    assert district_a["issues_requires_work_pct"] == 67
    assert district_a["issues_cohort_closed_pct"] == 50
    assert district_b["issues_found"] == 1
    assert district_b["issues_snapshot_total"] == 1
    assert district_b["issues_requires_work_pct"] == 0
    assert district_b["issues_cohort_closed_pct"] == 100
    assert payload["totals"]["issues_requires_work_pct"] == 50
    assert payload["totals"]["issues_cohort_closed_pct"] == 67


@pytest.mark.asyncio
async def test_dashboard_uses_completed_at_and_computed_overdue(client, admin_headers):
    me = await client.get("/api/v1/auth/me", headers=admin_headers)
    user_id = me.json()["id"]
    district_id, courtyard_id, site_id, inspection_id = [str(uuid.uuid4()) for _ in range(4)]
    _exec(
        "INSERT INTO districts(id,name,code) VALUES (%(d)s,%(n)s,%(c)s);"
        "INSERT INTO courtyards(id,district_id,name) VALUES (%(y)s,%(d)s,%(yn)s);"
        "INSERT INTO sites(id,courtyard_id,type,area_m2,geometry,is_active) VALUES "
        "(%(s)s,%(y)s,'Детская площадка',100,ST_GeomFromText("
        "'POLYGON((39 55,39.01 55,39.01 55.01,39 55.01,39 55))',4326),true);"
        "INSERT INTO inspections(id,site_id,inspector_id,type,status,created_at,completed_at) VALUES "
        "(%(i)s,%(s)s,%(u)s,'regular','completed','2025-01-01T00:00:00Z','2026-08-19T10:00:00Z');"
        "INSERT INTO issues(inspection_id,site_id,category_id,title,status,due_date,created_by,created_at) VALUES "
        "(%(i)s,%(s)s,(SELECT id FROM issue_categories WHERE name='Прочее'),'Просрочка',"
        "'open','2026-08-18',%(u)s,'2026-08-19T11:00:00Z')",
        {"d": district_id, "n": f"Stats {district_id[:8]}", "c": district_id[:8],
         "y": courtyard_id, "yn": "Двор stats", "s": site_id,
         "i": inspection_id, "u": user_id},
    )
    response = await client.get(
        "/api/v1/stats/dashboard",
        params={"date_from": "2026-08-19", "date_to": "2026-08-19", "district_id": district_id},
        headers=admin_headers,
    )
    assert response.status_code == 200, response.text
    row = response.json()["districts"][0]
    assert row["inspections_total"] == 1
    # An issue without a defect answer is still not a green inspection.
    assert row["inspections_green"] == 0
    assert row["inspections_with_defects"] == 1
    assert row["issues_found"] == 1
    assert row["issues_overdue"] == 1

    # Current inventory coverage excludes a deactivated site, while event
    # metrics still retain the completed inspection for the selected period.
    _exec("UPDATE sites SET is_active=false WHERE id=%(s)s", {"s": site_id})
    inactive_response = await client.get(
        "/api/v1/stats/dashboard",
        params={"date_from": "2026-08-19", "date_to": "2026-08-19", "district_id": district_id},
        headers=admin_headers,
    )
    inactive_row = inactive_response.json()["districts"][0]
    assert inactive_row["total_sites"] == 0
    assert inactive_row["sites_inspected"] == 0
    assert inactive_row["inspections_total"] == 1


@pytest.mark.asyncio
async def test_dashboard_counts_fix_and_closure_events_for_the_selected_period(client, admin_headers):
    me = await client.get("/api/v1/auth/me", headers=admin_headers)
    user_id = me.json()["id"]
    district_id, courtyard_id, site_id, inspection_id, issue_id = [str(uuid.uuid4()) for _ in range(5)]
    _exec(
        "INSERT INTO districts(id,name,code) VALUES (%(d)s,%(n)s,%(c)s);"
        "INSERT INTO courtyards(id,district_id,name) VALUES (%(y)s,%(d)s,%(yn)s);"
        "INSERT INTO sites(id,courtyard_id,type,area_m2,geometry,is_active) VALUES "
        "(%(s)s,%(y)s,'Детская площадка',100,ST_GeomFromText("
        "'POLYGON((40 55,40.01 55,40.01 55.01,40 55.01,40 55))',4326),true);"
        "INSERT INTO inspections(id,site_id,inspector_id,type,status,created_at,completed_at) VALUES "
        "(%(i)s,%(s)s,%(u)s,'regular','completed','2026-08-18T08:00:00Z','2026-08-18T09:00:00Z');"
        "INSERT INTO issues(id,inspection_id,site_id,category_id,title,status,created_by,created_at) VALUES "
        "(%(q)s,%(i)s,%(s)s,(SELECT id FROM issue_categories WHERE name='Прочее'),'Исправление вне дня создания',"
        "'closed',%(u)s,'2026-08-18T10:00:00Z');"
        "INSERT INTO issue_status_history(id,issue_id,old_status,new_status,changed_by,created_at) VALUES "
        "(%(h1)s,%(q)s,'in_work','fixed',%(u)s,'2026-08-19T10:00:00Z'),"
        "(%(h2)s,%(q)s,'fixed','closed',%(u)s,'2026-08-19T11:00:00Z');",
        {"d": district_id, "n": f"Events {district_id[:8]}", "c": district_id[:8],
         "y": courtyard_id, "yn": "Двор events", "s": site_id, "i": inspection_id,
         "q": issue_id, "h1": str(uuid.uuid4()), "h2": str(uuid.uuid4()), "u": user_id},
    )
    response = await client.get(
        "/api/v1/stats/dashboard",
        params={"date_from": "2026-08-19", "date_to": "2026-08-19", "district_id": district_id},
        headers=admin_headers,
    )
    assert response.status_code == 200, response.text
    row = response.json()["districts"][0]
    # The issue was created yesterday, but its repair and closure both happened today.
    assert row["issues_found"] == 0
    assert row["issues_fixed_events"] == 1
    assert row["issues_closed_events"] == 1
    assert row["issues_revision_events"] == 0
    assert row["issues_pending_final_current"] == 0
    assert row["issues_requires_work_current"] == 0


@pytest.mark.asyncio
async def test_admin_cannot_use_reviewer_bulk_accept(client, admin_headers):
    response = await client.post(
        "/api/v1/inspections/bulk-accept",
        json={"ids": [str(uuid.uuid4())]},
        headers=admin_headers,
    )
    assert response.status_code == 403
    assert response.json()["detail"] == "Массовая приёмка — только для проверяющего"


@pytest.mark.asyncio
async def test_dashboard_query_count_does_not_grow_with_districts():
    from app.database import async_session, engine

    statements = []

    def record_statement(*args):
        statements.append(args[2])

    event.listen(engine.sync_engine, "before_cursor_execute", record_statement)
    try:
        async with async_session() as db:
            user = SimpleNamespace(role="admin", district_id=None)
            filters = build_filter(user, date(2026, 8, 18), date(2026, 8, 20), None)
            await StatisticsService(db, filters).dashboard()
    finally:
        event.remove(engine.sync_engine, "before_cursor_execute", record_statement)
        await engine.dispose()

    assert len(statements) == 7


@pytest.mark.asyncio
async def test_stats_sections_requires_district_id(client, admin_headers):
    response = await client.get(
        "/api/v1/stats/sections",
        params={"date_from": "2026-08-18", "date_to": "2026-08-20"},
        headers=admin_headers,
    )
    assert response.status_code == 422


@pytest.mark.asyncio
async def test_stats_sections_breaks_down_by_courtyard_section(client, admin_headers):
    me = await client.get("/api/v1/auth/me", headers=admin_headers)
    user_id = me.json()["id"]
    district_id = str(uuid.uuid4())
    courtyard_a, courtyard_b, courtyard_c = [str(uuid.uuid4()) for _ in range(3)]
    site_a, site_b, site_c = [str(uuid.uuid4()) for _ in range(3)]
    inspection_a, inspection_b, inspection_c = [str(uuid.uuid4()) for _ in range(3)]
    _exec(
        "INSERT INTO districts(id,name,code) VALUES (%(d)s,%(n)s,%(c)s);"
        "INSERT INTO courtyards(id,district_id,name,section) VALUES "
        "(%(ya)s,%(d)s,'Двор А','Участок 1'),"
        "(%(yb)s,%(d)s,'Двор Б','Участок 2'),"
        "(%(yc)s,%(d)s,'Двор В',NULL);"
        "INSERT INTO sites(id,courtyard_id,type,area_m2,geometry,is_active) VALUES "
        "(%(sa)s,%(ya)s,'Детская площадка',100,ST_GeomFromText("
        "'POLYGON((40 55,40.01 55,40.01 55.01,40 55.01,40 55))',4326),true),"
        "(%(sb)s,%(yb)s,'Детская площадка',100,ST_GeomFromText("
        "'POLYGON((41 55,41.01 55,41.01 55.01,41 55.01,41 55))',4326),true),"
        "(%(sc)s,%(yc)s,'Детская площадка',100,ST_GeomFromText("
        "'POLYGON((42 55,42.01 55,42.01 55.01,42 55.01,42 55))',4326),true);"
        "INSERT INTO inspections(id,site_id,inspector_id,type,status,created_at,completed_at) VALUES "
        "(%(ia)s,%(sa)s,%(u)s,'regular','completed','2026-08-19T08:00:00Z','2026-08-19T09:00:00Z'),"
        "(%(ib)s,%(sb)s,%(u)s,'regular','completed','2026-08-19T08:00:00Z','2026-08-19T09:00:00Z'),"
        "(%(ic)s,%(sc)s,%(u)s,'regular','completed','2026-08-19T08:00:00Z','2026-08-19T09:00:00Z');",
        {"d": district_id, "n": f"Sections {district_id[:8]}", "c": district_id[:8],
         "ya": courtyard_a, "yb": courtyard_b, "yc": courtyard_c,
         "sa": site_a, "sb": site_b, "sc": site_c,
         "ia": inspection_a, "ib": inspection_b, "ic": inspection_c, "u": user_id},
    )

    response = await client.get(
        "/api/v1/stats/sections",
        params={"date_from": "2026-08-19", "date_to": "2026-08-19", "district_id": district_id},
        headers=admin_headers,
    )
    assert response.status_code == 200, response.text
    payload = response.json()
    assert payload["district_id"] == district_id
    by_section = {row["section"]: row for row in payload["sections"]}
    assert set(by_section) == {"Участок 1", "Участок 2", "Без участка"}
    assert by_section["Участок 1"]["total_sites"] == 1
    assert by_section["Участок 1"]["sites_inspected"] == 1
    assert by_section["Участок 2"]["total_sites"] == 1
    assert by_section["Без участка"]["total_sites"] == 1
    assert payload["totals"]["total_sites"] == 3
    assert payload["totals"]["sites_inspected"] == 3
    assert payload["totals"]["inspections_total"] == 3

    # На окружной штаб этот разрез не подаётся — /dashboard не меняется и
    # видит те же три площадки одной строкой по району.
    dashboard = await client.get(
        "/api/v1/stats/dashboard",
        params={"date_from": "2026-08-19", "date_to": "2026-08-19", "district_id": district_id},
        headers=admin_headers,
    )
    assert dashboard.status_code == 200, dashboard.text
    assert dashboard.json()["districts"][0]["total_sites"] == 3
