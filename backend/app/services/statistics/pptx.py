"""Deterministic two-slide headquarters presentation renderer."""

from io import BytesIO

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

from app.schemas import StatsCategoriesOut, StatsDashboardOut
from app.services.timezone import MSK

FONT = "Century Gothic"
RED = "E30613"
FRAME = "D9D9D9"
HEADER = "B4C7E7"
TOTAL = "595959"


def percentage_color(value: int | None) -> str:
    if value is None:
        return "D9E2F3"
    if value >= 90:
        return "63BE7B"
    if value >= 75:
        return "A9D18E"
    if value >= 60:
        return "FFD966"
    if value >= 40:
        return "F9CB9C"
    if value >= 20:
        return "F4B183"
    return "E06666"


def metric_label(numerator: int, denominator: int, percentage: int | None) -> str:
    if percentage is None or denominator == 0:
        return "—"
    return f"{numerator} из {denominator} · {percentage}%"


def _fill(cell, hex_color: str):
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor.from_string(hex_color)


def _style_cell(cell, *, size=8, bold=False, color="222222", align=PP_ALIGN.CENTER):
    cell.margin_left = cell.margin_right = Inches(0.03)
    cell.margin_top = cell.margin_bottom = Inches(0.02)
    for paragraph in cell.text_frame.paragraphs:
        paragraph.alignment = align
        for run in paragraph.runs:
            run.font.name = FONT
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = RGBColor.from_string(color)


def _shape(slide, shape_type, left, top, width, height, color, *, line_color=None):
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor.from_string(color)
    shape.line.color.rgb = RGBColor.from_string(line_color or color)
    return shape


def _add_chrome(slide, section: str):
    _shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 13.333, 0.63, "EFEFEF")
    _shape(slide, MSO_SHAPE.RECTANGLE, 0.0, 7.43, 13.333, 0.07, RED)
    _shape(slide, MSO_SHAPE.RECTANGLE, 12.76, 0.08, 0.45, 0.46, "FFFFFF")
    _shape(slide, MSO_SHAPE.RECTANGLE, 12.73, 0.08, 0.02, 0.46, RED)
    badge = slide.shapes.add_textbox(Inches(12.78), Inches(0.17), Inches(0.4), Inches(0.2))
    badge.text_frame.paragraphs[0].text = section
    badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    for run in badge.text_frame.paragraphs[0].runs:
        run.font.name, run.font.size, run.font.bold = FONT, Pt(10), True
    mark = _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 0.28, 0.13, 0.32, 0.32, RED)
    mark.text_frame.paragraphs[0].text = "ЖКХ"
    for run in mark.text_frame.paragraphs[0].runs:
        run.font.name, run.font.size, run.font.bold, run.font.color.rgb = FONT, Pt(5), True, RGBColor(255, 255, 255)
    unit = slide.shapes.add_textbox(Inches(0.68), Inches(0.12), Inches(3.4), Inches(0.35))
    unit.text_frame.paragraphs[0].text = "УПРАВЛЕНИЕ ЖИЛИЩНО-КОММУНАЛЬНОГО ХОЗЯЙСТВА"
    for run in unit.text_frame.paragraphs[0].runs:
        run.font.name, run.font.size, run.font.bold = FONT, Pt(7), True


def _add_title(slide, text: str, top=0.72):
    box = slide.shapes.add_textbox(Inches(3.55), Inches(top), Inches(8.7), Inches(0.34))
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = text
    paragraph.alignment = PP_ALIGN.CENTER
    run = paragraph.runs[0]
    run.font.name = FONT
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = RGBColor.from_string("4D4D4D")


def _add_metadata(slide, dashboard: StatsDashboardOut):
    generated = dashboard.generated_at.astimezone(MSK).strftime("%d.%m.%Y %H:%M")
    box = slide.shapes.add_textbox(Inches(9.1), Inches(7.08), Inches(3.35), Inches(0.16))
    paragraph = box.text_frame.paragraphs[0]
    paragraph.text = f"Методика v2 · МСК (UTC+3) · сформировано {generated}"
    paragraph.alignment = PP_ALIGN.RIGHT
    for run in paragraph.runs:
        run.font.name, run.font.size = FONT, Pt(6)
        run.font.color.rgb = RGBColor.from_string("777777")


def render_shtab(dashboard: StatsDashboardOut, categories: StatsCategoriesOut) -> BytesIO:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    _add_chrome(slide, "1.1")
    _add_title(slide, "Обходы детских и спортивных площадок САО — итоги периода")

    rows = sorted(
        dashboard.districts,
        key=lambda r: (-(r.clean_sites_pct if r.clean_sites_pct is not None else -1), -r.coverage_pct, r.district_name),
    )
    headers = ["№", "Район", "Охват", "Чистые площадки", "Площадки с нарушениями",
               "Без нарушений", "С нарушениями"]
    table = slide.shapes.add_table(
        len(rows) + 2, len(headers), Inches(0.35), Inches(1.12), Inches(12.63), Inches(5.0)
    ).table
    widths = [0.35, 2.15, 1.4, 2.0, 2.05, 1.3, 1.2]
    for column, width in zip(table.columns, widths):
        column.width = Inches(width)
    for index, header in enumerate(headers):
        table.cell(0, index).text = header
        _fill(table.cell(0, index), HEADER)
        _style_cell(table.cell(0, index), bold=True, color="FFFFFF")
    for row_index, row in enumerate(rows, start=1):
        clean = (
            f"{row.sites_latest_clean} из {row.sites_inspected} · {row.clean_sites_pct}%"
            if row.clean_sites_pct is not None else "Нет обходов"
        )
        defects = (
            f"{row.sites_latest_with_defects} из {row.sites_inspected} · {row.defect_sites_pct}%"
            if row.defect_sites_pct is not None else "Нет обходов"
        )
        values = [
            row_index, row.district_name,
            f"{row.sites_inspected} из {row.total_sites} · {row.coverage_pct}%",
            clean, defects, row.inspections_green, row.inspections_with_defects,
        ]
        for col, value in enumerate(values):
            table.cell(row_index, col).text = str(value)
            _style_cell(table.cell(row_index, col), align=PP_ALIGN.LEFT if col == 1 else PP_ALIGN.CENTER)
        _fill(table.cell(row_index, 2), percentage_color(row.coverage_pct))
        _fill(table.cell(row_index, 3), percentage_color(row.clean_sites_pct))
        _fill(
            table.cell(row_index, 4),
            percentage_color(100 - row.defect_sites_pct if row.defect_sites_pct is not None else None),
        )
    total_index = len(rows) + 1
    total = dashboard.totals
    total_values = [
        "", "ИТОГО", f"{total.sites_inspected} из {total.total_sites} · {total.coverage_pct}%",
        (
            f"{total.sites_latest_clean} из {total.sites_inspected} · {total.clean_sites_pct}%"
            if total.clean_sites_pct is not None else "Нет обходов"
        ),
        (
            f"{total.sites_latest_with_defects} из {total.sites_inspected} · {total.defect_sites_pct}%"
            if total.defect_sites_pct is not None else "Нет обходов"
        ),
        total.inspections_green, total.inspections_with_defects,
    ]
    for col, value in enumerate(total_values):
        table.cell(total_index, col).text = str(value)
        _fill(table.cell(total_index, col), TOTAL)
        _style_cell(table.cell(total_index, col), bold=True)

    p = dashboard.period
    summary = (
        f"Период: {p.date_from:%d.%m.%Y}–{p.date_to:%d.%m.%Y} МСК (UTC+3). "
        f"Обойдены {total.sites_inspected} из {total.total_sites} площадок ({total.coverage_pct}%); "
        f"чистые по последнему обходу в периоде — {total.sites_latest_clean} из "
        f"{total.sites_inspected} ({total.clean_sites_pct if total.clean_sites_pct is not None else '—'}%). "
        f"Площадки с нарушениями — {total.sites_latest_with_defects} из {total.sites_inspected} "
        f"({total.defect_sites_pct if total.defect_sites_pct is not None else '—'}%)."
    )
    _shape(slide, MSO_SHAPE.RECTANGLE, 0.35, 6.32, 0.02, 0.7, RED)
    box = slide.shapes.add_textbox(Inches(0.48), Inches(6.34), Inches(12.0), Inches(0.6))
    box.text_frame.paragraphs[0].text = summary
    for run in box.text_frame.paragraphs[0].runs:
        run.font.name, run.font.size = FONT, Pt(11)
    _add_metadata(slide, dashboard)

    slide2 = prs.slides.add_slide(blank)
    _add_chrome(slide2, "1.2")
    _add_title(slide2, "Устранение замечаний")
    for left, width, text in (
        (0.95, 3.6, "Поток за период"),
        (5.2, 2.25, "Результат по замечаниям периода"),
        (7.75, 4.6, "Состояние на конец периода"),
    ):
        box = slide2.shapes.add_textbox(Inches(left), Inches(1.0), Inches(width), Inches(0.22))
        paragraph = box.text_frame.paragraphs[0]
        paragraph.text = text
        paragraph.alignment = PP_ALIGN.CENTER
        for run in paragraph.runs:
            run.font.name, run.font.size, run.font.bold = FONT, Pt(9), True
            run.font.color.rgb = RGBColor.from_string("4D4D4D")

    remediation_table = slide2.shapes.add_table(
        len(rows) + 2, 11, Inches(0.35), Inches(1.25), Inches(12.63), Inches(4.95)
    ).table
    remediation_headers = [
        "№", "Район", "Выявлено", "На финальной проверке", "Исправлено", "Доработка",
        "Устранено из выявленных", "На проверке", "Требуют устранения", "Просрочено",
        "Доля требующих устранения",
    ]
    widths = [0.35, 1.85, 0.75, 1.05, 0.8, 0.8, 1.75, 0.85, 1.05, 0.8, 1.6]
    for column, width in zip(remediation_table.columns, widths):
        column.width = Inches(width)
    for index, header in enumerate(remediation_headers):
        remediation_table.cell(0, index).text = header
        _fill(remediation_table.cell(0, index), HEADER)
        _style_cell(remediation_table.cell(0, index), bold=True, color="FFFFFF")
    for row_index, row in enumerate(rows, start=1):
        values = [
            row_index,
            row.district_name,
            row.issues_found,
            row.issues_fixed_events,
            row.issues_closed_events,
            row.issues_revision_events,
            metric_label(row.issues_cohort_closed_as_of, row.issues_found, row.issues_cohort_closed_pct),
            row.issues_pending_final_current,
            row.issues_requires_work_current,
            row.issues_overdue_current,
            metric_label(
                row.issues_requires_work_current, row.issues_snapshot_total, row.issues_requires_work_pct,
            ),
        ]
        for col, value in enumerate(values):
            remediation_table.cell(row_index, col).text = str(value)
            _style_cell(
                remediation_table.cell(row_index, col),
                align=PP_ALIGN.LEFT if col == 1 else PP_ALIGN.CENTER,
            )
        _fill(remediation_table.cell(row_index, 6), percentage_color(row.issues_cohort_closed_pct))
        _fill(
            remediation_table.cell(row_index, 10),
            percentage_color(
                100 - row.issues_requires_work_pct
                if row.issues_requires_work_pct is not None else None
            ),
        )

    remediation_total = [
        "", "ИТОГО", total.issues_found, total.issues_fixed_events, total.issues_closed_events,
        total.issues_revision_events,
        metric_label(
            total.issues_cohort_closed_as_of, total.issues_found, total.issues_cohort_closed_pct,
        ),
        total.issues_pending_final_current,
        total.issues_requires_work_current,
        total.issues_overdue_current,
        metric_label(
            total.issues_requires_work_current, total.issues_snapshot_total,
            total.issues_requires_work_pct,
        ),
    ]
    remediation_total_index = len(rows) + 1
    for col, value in enumerate(remediation_total):
        remediation_table.cell(remediation_total_index, col).text = str(value)
        _fill(remediation_table.cell(remediation_total_index, col), TOTAL)
        _style_cell(remediation_table.cell(remediation_total_index, col), bold=True)

    summary2 = (
        f"Период: {p.date_from:%d.%m.%Y}–{p.date_to:%d.%m.%Y} МСК (UTC+3). "
        f"Устранено из выявленных замечаний периода — "
        f"{metric_label(total.issues_cohort_closed_as_of, total.issues_found, total.issues_cohort_closed_pct)}. "
        f"Доля требующих устранения на конец периода — "
        f"{metric_label(total.issues_requires_work_current, total.issues_snapshot_total, total.issues_requires_work_pct)}."
    )
    _shape(slide2, MSO_SHAPE.RECTANGLE, 0.35, 6.42, 0.02, 0.52, RED)
    note = slide2.shapes.add_textbox(Inches(0.48), Inches(6.45), Inches(12.0), Inches(0.42))
    note.text_frame.paragraphs[0].text = summary2
    for run in note.text_frame.paragraphs[0].runs:
        run.font.name, run.font.size = FONT, Pt(10)
    _add_metadata(slide2, dashboard)

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output
